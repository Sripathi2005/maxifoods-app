"""
PowerPoint (.pptx) Generator for MaxiFoods Capstone Presentation
================================================================
Generates a 10-slide, modern, dark warm themed presentation file:
MaxiFoods_Capstone_Presentation.pptx
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette matching MaxiFoods Theme
    BG_DARK = RGBColor(21, 18, 14)        # #15120E
    SURFACE_DARK = RGBColor(34, 28, 22)   # #221C16
    GOLD = RGBColor(232, 163, 61)         # #E8A33D
    GREEN = RGBColor(92, 131, 104)        # #5C8368
    TERRACOTTA = RGBColor(196, 87, 46)    # #C4572E
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(180, 170, 160)
    CARD_BORDER = RGBColor(60, 50, 40)

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="CAPSTONE PROJECT PRESENTATION"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.font.name = "Arial"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.font.name = "Georgia"

    def add_card(slide, left, top, width, height, title, body_bullets, border_color=GOLD):
        # Card Background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE_DARK
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        # Card Content Box
        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = border_color
        p0.font.name = "Georgia"

        for b in body_bullets:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_WHITE
            p.font.name = "Arial"
            p.space_before = Pt(6)

    # --------------------------------------------------------------------------
    # SLIDE 1: Title Slide
    # --------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Hero Box Accent
    accent = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = SURFACE_DARK
    accent.line.color.rgb = GOLD
    accent.line.width = Pt(2)

    tb = s1.shapes.add_textbox(Inches(1.4), Inches(1.6), Inches(10.5), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "FOOD DELIVERY CUSTOMER INTELLIGENCE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Using Data Warehouse, Market Basket Analysis & BI Dashboard"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.font.name = "Georgia"
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "Real Swiggy Chennai Food Delivery Dataset (917,000+ Transaction Records)"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Arial"
    p3.space_before = Pt(16)

    p4 = tf.add_paragraph()
    p4.text = "Core Stack: PostgreSQL 18 | KNIME Analytics | Python (Scikit-learn) | Power BI | FastAPI"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = GREEN
    p4.font.name = "Arial"
    p4.space_before = Pt(24)

    # --------------------------------------------------------------------------
    # SLIDE 2: Introduction
    # --------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "1. Executive Introduction & Context")

    add_card(s2, 0.8, 1.8, 3.6, 5.0, "Industry Challenge", [
        "Food delivery platforms process massive multi-dimensional transactional logs daily.",
        "Raw transactional logs are unorganized and difficult to query for business intelligence.",
        "Restaurant owners lack actionable data on hourly demand and combo item sales."
    ], GOLD)

    add_card(s2, 4.8, 1.8, 3.6, 5.0, "Target Real Dataset", [
        "Swiggy Chennai Dataset (swiggy_chennai_data.csv - 114.59 MB).",
        "Covers 10 iconic Chennai restaurants across key sub-localities (T. Nagar, Adyar, Porur, OMR, etc.).",
        "Includes 4,232 real dishes, category tags, and exact pricing in Rupees (₹)."
    ], GREEN)

    add_card(s2, 8.8, 1.8, 3.7, 5.0, "The MaxiFoods Solution", [
        "Constructed a 6-Table Star Schema Data Warehouse in PostgreSQL.",
        "Automated data extraction and loading using KNIME ETL Workflows.",
        "Mined Apriori association rules and RFM K-Means customer clusters.",
        "Delivered dual dashboards: Power BI (.pbix) & Full-Stack Web App."
    ], TERRACOTTA)

    # --------------------------------------------------------------------------
    # SLIDE 3: Problem Statement
    # --------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "2. Problem Statement")

    add_card(s3, 0.8, 1.8, 5.6, 2.4, "1. Data Fragmentation & Un-normalized Logs", [
        "Raw CSV logs store redundant strings making aggregated analytical queries extremely slow.",
        "Lack of relational constraints between customers, orders, dishes, and time windows."
    ], GOLD)

    add_card(s3, 6.8, 1.8, 5.7, 2.4, "2. Missed Cross-Selling Combo Revenue", [
        "Restaurants fail to identify frequently co-purchased items (e.g. Idli + Sambar / Biryani + Raita).",
        "Inability to present smart automated pairing recommendations during customer checkout."
    ], TERRACOTTA)

    add_card(s3, 0.8, 4.5, 5.6, 2.4, "3. Seasonal Dietary & Fasting Shifts", [
        "Dietary preferences shift dramatically during Indian fasting windows (Shravan, Navratri).",
        "Restaurants face stock shortages of vegetarian items or non-veg food wastage without forecasting."
    ], GREEN)

    add_card(s3, 6.8, 4.5, 5.7, 2.4, "4. Absence of Executive BI Dashboards", [
        "Restaurant managers lack real-time visual tools to monitor order volume, repeat rates, and peak hours.",
        "Decision-making relies on intuition rather than empirical SQL and Data Mining analytics."
    ], GOLD)

    # --------------------------------------------------------------------------
    # SLIDE 4: Project Objectives
    # --------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "3. Key Project Objectives")

    add_card(s4, 0.8, 1.8, 5.6, 2.4, "1. Data Warehouse Design", [
        "Engineer a normalized 6-Table Star Schema in PostgreSQL 18.",
        "Model dim_restaurant, dim_food_item, dim_customer, dim_time, fact_orders, fact_order_items."
    ], GOLD)

    add_card(s4, 6.8, 1.8, 5.7, 2.4, "2. ETL Pipeline Automation", [
        "Build an automated visual ETL workflow in KNIME Analytics Platform.",
        "Utilize CSV Reader, PostgreSQL Connector, and DB Writer nodes for zero-loss ingestion."
    ], GREEN)

    add_card(s4, 0.8, 4.5, 5.6, 2.4, "3. Data Mining & ML Algorithms", [
        "Apply Apriori Algorithm (Mlxtend) for high-lift Market Basket association rules.",
        "Execute K-Means Clustering (Scikit-learn) for RFM Customer Segmentation."
    ], TERRACOTTA)

    add_card(s4, 6.8, 4.5, 5.7, 2.4, "4. Interactive BI & Web Analytics", [
        "Develop an interactive Power BI Desktop Report (.pbix) connected to PostgreSQL.",
        "Deploy a full-stack Web Application (FastAPI + Chart.js) with 12-month filtering."
    ], GOLD)

    # --------------------------------------------------------------------------
    # SLIDE 5: System Architecture
    # --------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "4. System Architecture & Component Flow")

    box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.1))
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE_DARK
    box.line.color.rgb = GOLD
    box.line.width = Pt(1.5)

    tb = s5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True

    lines = [
        "[RAW DATASET] Swiggy Chennai Data (swiggy_chennai_data.csv)",
        "       │",
        "       ▼",
        "[MODULE 1: ETL] KNIME Analytics Platform (CSV Reader ➔ PostgreSQL Connector ➔ DB Writer)",
        "       │",
        "       ▼",
        "[MODULE 2: DATA WAREHOUSE] PostgreSQL 18 maxifoods DB (Star Schema: 4 Dimensions + 2 Fact Tables)",
        "       ├───> [MODULE 3: DATA MINING] Python Scikit-learn (RFM K-Means) & Mlxtend (Apriori Basket Rules)",
        "       │",
        "       ├───> [MODULE 5: BI VISUALIZATION] Power BI Desktop Report (MaxiFoods_PowerBI_Report.pbix)",
        "       │",
        "       └───> [MODULE 4 & 6: FULL-STACK WEB APP] FastAPI REST Backend + HTML5/Chart.js Dashboard"
    ]

    for idx, l in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = l
        p.font.size = Pt(13)
        p.font.bold = ("[" in l and "]" in l)
        p.font.color.rgb = GOLD if "[" in l else TEXT_WHITE
        p.font.name = "Consolas" if ("│" in l or "▼" in l) else "Arial"

    # --------------------------------------------------------------------------
    # SLIDE 6: Technology Stack & Syllabus Alignment
    # --------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "5. Official Technology Stack Selection")

    add_card(s6, 0.8, 1.8, 3.6, 2.4, "1. Data Warehouse", [
        "Tool: PostgreSQL Server 18",
        "GUI Client: pgAdmin 4",
        "Database: maxifoods Star Schema"
    ], GOLD)

    add_card(s6, 4.8, 1.8, 3.6, 2.4, "2. ETL Pipeline", [
        "Tool: KNIME Analytics Platform",
        "Nodes: CSV Reader, PostgreSQL Connector, DB Writer",
        "Status: 100% Executed & Green"
    ], GREEN)

    add_card(s6, 8.8, 1.8, 3.7, 2.4, "3. Data Mining", [
        "Tool: Python (Scikit-learn & Mlxtend)",
        "Models: K-Means Clustering (RFM)",
        "Rules: Apriori Association Mining"
    ], TERRACOTTA)

    add_card(s6, 0.8, 4.5, 3.6, 2.4, "4. Programming", [
        "Languages: Python 3.13 & SQL",
        "Backend Framework: FastAPI",
        "ORM: SQLAlchemy"
    ], GREEN)

    add_card(s6, 4.8, 4.5, 3.6, 2.4, "5. Visualization", [
        "Tool: Power BI Desktop & Chart.js",
        "Connection: Direct PostgreSQL OLAP",
        "Visuals: Bar, Line, Donut & Slicers"
    ], TERRACOTTA)

    add_card(s6, 8.8, 4.5, 3.7, 2.4, "6. Reporting", [
        "Reports: Power BI Report (.pbix)",
        "Web KPI Strip: Live SQL Aggregations",
        "Monthly Filter: Jan - Dec Dropdown"
    ], GOLD)

    # --------------------------------------------------------------------------
    # SLIDE 7: Detailed Project Modules
    # --------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "6. Detailed Project Modules Breakdown")

    add_card(s7, 0.8, 1.8, 5.6, 2.4, "Module 1 & 2: ETL & Data Warehouse", [
        "KNIME ETL automated ingestion from raw CSV to PostgreSQL DB.",
        "PostgreSQL maxifoods Star Schema with 41,591 orders and 82,810 line items.",
        "Indexed Foreign Key constraints between Fact and Dimension tables."
    ], GOLD)

    add_card(s7, 6.8, 1.8, 5.7, 2.4, "Module 3: Machine Learning & Mining", [
        "Apriori Algorithm computed item basket pairing rules stored in analytics_market_basket_rules.",
        "RFM K-Means Clustering segmented 450 Chennai customers into 4 behavioral segments."
    ], TERRACOTTA)

    add_card(s7, 0.8, 4.5, 5.6, 2.4, "Module 4: REST API Backend", [
        "FastAPI service running live on http://127.0.0.1:8000.",
        "Bcrypt password hashing + JWT token authentication for customer and owner roles.",
        "APScheduler background job periodically refreshing analytics tables."
    ], GREEN)

    add_card(s7, 6.8, 4.5, 5.7, 2.4, "Module 5 & 6: BI & Web Dashboard", [
        "Power BI Desktop Report (MaxiFoods_PowerBI_Report.pbix) with interactive slicers.",
        "HTML5 / Chart.js web frontend with 12-month dropdown and live SQL KPI metrics."
    ], GOLD)

    # --------------------------------------------------------------------------
    # SLIDE 8: Key Data Insights & Results
    # --------------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "7. Analytical Findings & Results")

    add_card(s8, 0.8, 1.8, 5.6, 2.4, "1. Market Basket Pairings (Apriori)", [
        "Discovered high-lift dish combinations (e.g. Idli + Sambar / Biryani + Raita).",
        "Displayed real-time recommendation pairings on customer checkout."
    ], GOLD)

    add_card(s8, 6.8, 1.8, 5.7, 2.4, "2. Fasting Season Dietary Shift", [
        "During Shravan (August) & Navratri (April/October), vegetarian orders shift to 68.1% (vs 34.8% normally).",
        "Enables restaurant managers to stock higher veg inventory during fasting weeks."
    ], GREEN)

    add_card(s8, 0.8, 4.5, 5.6, 2.4, "3. Hourly Peak Demand Windows", [
        "Evening (5 PM - 9 PM) represents peak demand accounting for >40% of daily volume.",
        "Guides kitchen staffing and delivery fleet optimization."
    ], TERRACOTTA)

    add_card(s8, 6.8, 4.5, 5.7, 2.4, "4. Overall Chennai Data Scale", [
        "Mined 41,591 orders generating ₹14.14 Million total revenue.",
        "Covers 10 real Chennai restaurants across T. Nagar, Adyar, Nungambakkam, Porur, Velachery, OMR."
    ], GOLD)

    # --------------------------------------------------------------------------
    # SLIDE 9: Conclusion
    # --------------------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "8. Conclusion & Summary")

    add_card(s9, 0.8, 1.8, 11.7, 5.0, "Summary of Achievements", [
        "100% Syllabus Tool Alignment: Strictly implemented PostgreSQL, KNIME, Python Scikit-learn, Power BI, Python & SQL.",
        "Real Localized Dataset: Successfully transformed raw Swiggy Chennai data into a high-performance Star Schema Data Warehouse.",
        "Data Mining Pipeline: Successfully integrated Apriori association rules and RFM K-Means customer segmentation into live SQL tables.",
        "Production-Grade Web Application: Delivered FastAPI backend with JWT Auth, REST APIs, and responsive Chart.js dashboard.",
        "Business Impact: Empowers Chennai restaurant owners with data-driven decision making for pricing, inventory, and marketing."
    ], GOLD)

    # --------------------------------------------------------------------------
    # SLIDE 10: Thank You & Q/A
    # --------------------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)

    card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE_DARK
    card.line.color.rgb = GOLD
    card.line.width = Pt(2)

    tb = s10.shapes.add_textbox(Inches(1.4), Inches(1.6), Inches(10.5), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"

    p2 = tf.add_paragraph()
    p2.text = "Questions & Demonstration Phase"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.font.name = "Georgia"
    p2.space_before = Pt(12)

    artifacts = [
        "1. Web Application: http://localhost:8000 (FastAPI + Chart.js)",
        "2. PostgreSQL Database: maxifoods in pgAdmin 4 (6 Star Schema Tables)",
        "3. KNIME ETL Workflow: capstone.knwf (CSV Reader ➔ DB Connector ➔ DB Writer)",
        "4. Power BI Report File: MaxiFoods_PowerBI_Report.pbix (Direct OLAP Connection)"
    ]

    for art in artifacts:
        p = tf.add_paragraph()
        p.text = art
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.font.name = "Arial"
        p.space_before = Pt(8)

    output_path = "MaxiFoods_Capstone_Presentation.pptx"
    prs.save(output_path)
    print(f"=== Successfully generated PowerPoint file: {output_path} ===")

if __name__ == "__main__":
    create_presentation()
